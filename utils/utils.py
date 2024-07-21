# from langchain_google_genai import ChatGoogleGenerativeAI
# from langchain_google_genai import (
#     ChatGoogleGenerativeAI,
#     HarmBlockThreshold,
#     HarmCategory,
# )
# from langchain_community.document_loaders import PyMuPDFLoader
# from langchain.chains.llm import LLMChain
import fitz
import io
from tqdm import tqdm
from transformers import AutoProcessor, AutoModelForCausalLM
import gc
import torch
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO


def pptx_to_images(pptx_path, scale=0.5):
    pages = []
    prs = Presentation(pptx_path)

    # Calculate scaled slide dimensions
    slide_width = int(prs.slide_width.pt * scale)
    slide_height = int(prs.slide_height.pt * scale)

    for i, slide in enumerate(prs.slides):
        # Create a blank image with white background
        slide_image = Image.new('RGB', (slide_width, slide_height), 'white')
        draw = ImageDraw.Draw(slide_image)

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # PIL expects the image in bytes format
                img = Image.open(BytesIO(shape.image.blob))
                img = img.resize((int(shape.width.pt * scale),
                                 int(shape.height.pt * scale)))
                slide_image.paste(
                    img, (int(shape.left.pt * scale), int(shape.top.pt * scale)))
            elif shape.has_text_frame:
                left = int(shape.left.pt * scale)
                top = int(shape.top.pt * scale)
                width = int(shape.width.pt * scale)
                height = int(shape.height.pt * scale)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Set the font size and type
                        font_size = int(run.font.size.pt *
                                        scale) if run.font.size else 12
                        # Use DejaVuSans font for Unicode support
                        font = ImageFont.truetype(
                            "../DejaVuSans-Bold.ttf", font_size)
                        draw.text((left, top), run.text,
                                  fill="black", font=font)

        pages.append(slide_image)

    return pages


def convert_pdf_to_images(pdf_path):

    if (pdf_path.endswith(".docx")):
        # do something
        pass
    elif (pdf_path.endswith(".pdf")):
        pdf_document = fitz.open(pdf_path)
    elif (pdf_path.endswith(".pptx")):
        return pptx_to_images(pdf_path)

    pages = []

    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)

        zoom_x = 2.0
        zoom_y = 2.0
        mat = fitz.Matrix(zoom_x, zoom_y)

        pix = page.get_pixmap(matrix=mat)

        # Convert pixmap to image without saving to disk
        image_bytes = pix.tobytes("png")

        pages.append(Image.open(io.BytesIO(image_bytes)))

    return pages


def images_and_summarize(file_path, chroma_client):
    # Convert PDF to images (assuming this function is defined elsewhere)
    pages = convert_pdf_to_images(file_path)

    # Check for GPU availability
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    # device = 'cuda'

    # Load model and processor once
    model = AutoModelForCausalLM.from_pretrained(
        "microsoft/Florence-2-base", trust_remote_code=True).to(device)
    processor = AutoProcessor.from_pretrained(
        "microsoft/Florence-2-base", trust_remote_code=True)

    page_summaries = []

    for page_num in tqdm(range(len(pages))):
        image = pages[page_num]

        prompt = "<MORE_DETAILED_CAPTION>"

        # Process inputs
        inputs = processor(text=prompt, images=image, return_tensors="pt")
        input_ids = inputs["input_ids"].to(device)
        pixel_values = inputs["pixel_values"].to(device)

        # Generate text
        generated_ids = model.generate(
            input_ids=input_ids,
            pixel_values=pixel_values,
            max_new_tokens=1024,
            num_beams=3,
            do_sample=False
        )

        # Decode generated text
        generated_text = processor.batch_decode(
            generated_ids, skip_special_tokens=False)[0]

        # Post-process generation
        parsed_answer = processor.post_process_generation(
            generated_text, task="<MORE_DETAILED_CAPTION>", image_size=(image.width, image.height))

        page_summaries.append(parsed_answer['<MORE_DETAILED_CAPTION>'])
        # Clear memory
        gc.collect()
        torch.cuda.empty_cache()

    # Store embeddings
    client = chroma_client
    collection = client.get_or_create_collection(
        name='godrej_summaries',
        metadata={'hnsw:space': 'cosine', "userId": 1}
    )

    for i, doc in enumerate(page_summaries):
        collection.add(documents=[doc], ids=[str(i) + file_path])

    return page_summaries
